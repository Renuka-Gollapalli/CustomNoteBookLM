from pptx import Presentation

def load_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        text = " ".join([
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ])
        if text:
            chunks.append({
                "text": text,
                "page": i + 1,
                "source": file_path,
                "source_type": "pptx"
            })
    return chunks
